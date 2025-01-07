# utils/chart_updater.py
from pptx import Presentation
from pptx.chart.data import CategoryChartData
import pandas as pd
import numpy as np
from typing import Dict, List
from config.constants import PARTY_PAIRS, PARTY_2023_SLIDE

class ChartUpdater:
    def __init__(self, output_path: str, language: str = 'tr'):
        self.output_path = output_path
        self.language = language
        self.prs = None
        
        # Complete translation mappings for both Turkish and English
        self.translations = {
            'tr': {
                'months': {
                    'Oca': 'Oca',
                    'Şub': 'Şub',
                    'Mar': 'Mar',
                    'Nis': 'Nis',
                    'May': 'May',
                    'Haz': 'Haz',
                    'Tem': 'Tem',
                    'Ağu': 'Ağu',
                    'Eyl': 'Eyl',
                    'Eki': 'Eki',
                    'Kas': 'Kas',
                    'Ara': 'Ara'
                },
                'parties': {
                    'AK Parti': 'AK Parti',
                    'CHP': 'CHP',
                    'DEM Parti': 'DEM Parti',
                    'İYİ Parti': 'İYİ Parti',
                    'MHP': 'MHP',
                    'Kararsız': 'Kararsız',
                    'Oy Kullanmam': 'Oy Kullanmam',
                    'Oy kullanmayacağım': 'Oy kullanmayacağım',
                    'Diğer': 'Diğer'
                },
                'education': {
                    'İlköğretim ve altı': 'İlköğretim ve altı',
                    'Lise': 'Lise',
                    'Yüksekokul ve üzeri': 'Yüksekokul ve üzeri'
                },
                'age': {
                    '18-34': '18-34',
                    '35-54': '35-54',
                    '55 ve üstü': '55 ve üstü'
                },
                'economy': {
                    'Çok kötü / Kötü': 'Çok kötü / Kötü',
                    'Ne iyi ne kötü': 'Ne iyi ne kötü',
                    'Çok İyi / İyi': 'Çok İyi / İyi',
                    'Çok Daha Kötü/Daha Kötü': 'Çok Daha Kötü/Daha Kötü',
                    'Değişmez': 'Değişmez',
                    'Çok Daha İyi/Daha İyi': 'Çok Daha İyi/Daha İyi'
                },
                'subsistence': {
                    'Gelirim giderimi karşılamadı.': 'Gelirim giderimi karşılamadı.',
                    'Gelirim giderimi ucu ucuna karşıladı.': 'Gelirim giderimi ucu ucuna karşıladı.',
                    'Gelirim giderlerimin üzerinde oldu.': 'Gelirim giderlerimin üzerinde oldu.',
                    'Gelirim giderlerimi fazlasıyla karşıladı.': 'Gelirim giderlerimi fazlasıyla karşıladı.'
                },
                'chart_titles': {
                    'Oy Oranı': 'Oy Oranı',
                    'Success Rate': 'Success Rate'
                },
                'politicians': {
                    'Recep Tayyip Erdoğan': 'Recep Tayyip Erdoğan',
                    'Özgür Özel': 'Özgür Özel',
                    'Devlet Bahçeli': 'Devlet Bahçeli',
                    'Ekrem İmamoğlu': 'Ekrem İmamoğlu',
                    'Mansur Yavaş': 'Mansur Yavaş',
                    'Fatih Erbakan': 'Fatih Erbakan',
                    'Muharrem İnce': 'Muharrem İnce',
                    'Erkan Baş': 'Erkan Baş',
                    'Ümit Özdağ': 'Ümit Özdağ',
                    'Müsavat Dervişoğlu': 'Müsavat Dervişoğlu',
                    'Tülay Hatimoğulları Oruç': 'Tülay Hatimoğulları Oruç',
                    'Yavuz Ağıralioğlu': 'Yavuz Ağıralioğlu',
                    'Mahmut Arıkan': 'Mahmut Arıkan'
                }
            },
            'en': {
                'months': {
                    'Oca': 'Jan',
                    'Şub': 'Feb',
                    'Mar': 'Mar',
                    'Nis': 'Apr',
                    'May': 'May',
                    'Haz': 'Jun',
                    'Tem': 'Jul',
                    'Ağu': 'Aug',
                    'Eyl': 'Sep',
                    'Eki': 'Oct',
                    'Kas': 'Nov',
                    'Ara': 'Dec'
                },
                'parties': {
                    'AK Parti': 'AK Party',
                    'CHP': 'CHP',
                    'DEM Parti': 'HDP/Green Left',
                    'İYİ Parti': 'İYİ Party',
                    'MHP': 'MHP',
                    'Kararsız': 'Undecided',
                    'Kararsızım': 'Undecided',
                    'Oy Kullanmam': '''Won't Vote''',
                    'Oy kullanmayacağım': '''Won't Vote''',
                    'Diğer': 'Other',
                    'Zafer Partisi': 'Victory Party',
                    'Yeniden Refah Partisi': 'New Welfare Party',
                    'Anahtar Parti': 'Key Party',
                },
                'education': {
                    'İlköğretim ve altı': 'Primary school and below',
                    'Lise': 'High School',
                    'Yüksekokul ve üzeri': 'University and Above'
                },
                'age': {
                    '18-34': '18-34',
                    '35-54': '35-54',
                    '55 ve üstü': '55 and above'
                },
                'economy': {
                    'Çok kötü / Kötü': 'Very bad / Bad',
                    'Ne iyi ne kötü': 'Neither bad nor good',
                    'Çok İyi / İyi': 'Very good / Good',
                    'Çok Daha Kötü/Daha Kötü': 'Much Worse/Worse',
                    'Değişmez': '''Won't Change''',
                    'Çok Daha İyi/Daha İyi': 'Much Better/Better'
                },
                'subsistence': {
                    'Gelirim giderimi karşılamadı.': 'My income did not meet my expenses.',
                    'Gelirim giderimi ucu ucuna karşıladı.': 'My income barely met my expenses.',
                    'Gelirim giderlerimin üzerinde oldu.': 'My income exceeded my expenses.',
                    'Gelirim giderlerimi fazlasıyla karşıladı.': 'My income was well above my expenses.'
                },
                'chart_titles': {
                    'Oy Oranı': 'Vote Share',
                    'Success Rate': 'Success Rate'
                },
                'politicians': {
                    'Recep Tayyip Erdoğan': 'Recep Tayyip Erdoğan',
                    'Özgür Özel': 'Özgür Özel',
                    'Devlet Bahçeli': 'Devlet Bahçeli',
                    'Ekrem İmamoğlu': 'Ekrem İmamoğlu',
                    'Mansur Yavaş': 'Mansur Yavaş',
                    'Fatih Erbakan': 'Fatih Erbakan',
                    'Muharrem İnce': 'Muharrem İnce',
                    'Erkan Baş': 'Erkan Baş',
                    'Ümit Özdağ': 'Ümit Özdağ',
                    'Müsavat Dervişoğlu': 'Müsavat Dervişoğlu',
                    'Tülay Hatimoğulları Oruç': 'Tülay Hatimoğulları Oruç',
                    'Yavuz Ağıralioğlu': 'Yavuz Ağıralioğlu',
                    'Mahmut Arıkan': 'Mahmut Arıkan'
                }
            }
        }

    def _translate_date(self, date_str: str) -> str:
        """Translate date format from Turkish to English (e.g., 'Oca.23' to 'Jan.23')"""
        if self.language == 'tr' or not date_str:
            return date_str
        
        try:
            month, year = date_str.split('.')
            translated_month = self.translations[self.language]['months'].get(month, month)
            return f"{translated_month}.{year}"
        except:
            return date_str

    def _translate(self, text: str, category: str = 'parties') -> str:
        """Translate text based on language setting"""
        # Always get translations for the current language
        translations = self.translations.get(self.language, {}).get(category, {})
        translated = translations.get(text, text)
        
        # For debugging
        if self.language == 'en' and category == 'politicians':
            print(f"Translating '{text}' to '{translated}' (category: {category})")
        
        return translated

    def _translate_list(self, items: list, category: str = 'parties') -> list:
        """Translate a list of items"""
        if category == 'dates':
            return [self._translate_date(item) for item in items]
        return [self._translate(item, category) for item in items]

    def _load_presentation(self):
        """Load presentation if not already loaded"""
        if self.prs is None:
            self.prs = Presentation(self.output_path)

    def _save_presentation(self):
        """Save presentation and close it"""
        if self.prs:
            try:
                print(f"Saving presentation to: {self.output_path}")
                self.prs.save(self.output_path)
                print("Presentation saved successfully")
            except Exception as e:
                print(f"Error saving presentation: {str(e)}")
                raise
            finally:
                self.prs = None  # Close the presentation regardless of success or failure

    def update_all_charts(self, party_data: Dict[str, float], historical_data: Dict[str, pd.DataFrame]):
        """Update all charts in one go to avoid multiple file operations"""
        print(f"\nUpdating charts for language: {self.language}")
        print(f"Using presentation file: {self.output_path}")
        
        self._load_presentation()
        
        try:
            # Update main party chart
            self._update_party_chart(15, party_data)
            
            # Update time series charts
            self._update_time_series_charts(historical_data['party_votes'])
            
            # Update education charts
            self._update_education_charts(historical_data)
            
            # Update age charts
            self._update_age_charts(historical_data)
            
            # Update 2023 party chart
            self._update_2023_party_chart(historical_data['party_votes_2023'])
            
            # Update economic charts
            self._update_econ_charts(historical_data)
            
            # Update politician success charts - don't remove the data
            current_success_data = historical_data.get('current_success')
            if current_success_data is not None:
                self._update_politician_success_charts(historical_data, current_success_data)
            
            # Update subsistence charts
            self._update_subsistence_charts(historical_data)
            
            # Save the presentation after all updates
            print(f"Saving presentation to: {self.output_path}")
            self._save_presentation()
            
            # Verify the file exists and has size
            import os
            if os.path.exists(self.output_path):
                print(f"File saved successfully. Size: {os.path.getsize(self.output_path)} bytes")
            else:
                print("WARNING: File not found after saving!")
                
        except Exception as e:
            print(f"Error updating charts: {str(e)}")
            # Make sure to save even if there's an error
            self._save_presentation()
            raise

    def _update_party_chart(self, slide_number: int, party_data: Dict[str, float]):
        """Update the party voting chart"""
        slide = self.prs.slides[slide_number - 1]
        
        # Convert dictionary to DataFrame
        df_data = pd.DataFrame([(k, v) for k, v in party_data.items()], 
                             columns=['Party', 'Percentage'])
        
        # Sort data
        sorted_data = self._prepare_sorted_data(df_data)
        
        # Prepare chart data with translations
        chart_data = CategoryChartData(number_format='0.0%')
        chart_data.categories = [self._translate(party, 'parties') for party in sorted_data['Party'].tolist()]
        chart_data.add_series(self._translate('Oy Oranı', 'chart_titles'), [x/100 for x in sorted_data['Percentage'].tolist()])
        
        # Update Chart 1
        for shape in slide.shapes:
            if shape.has_chart and shape.name == "Chart 1":
                shape.chart.replace_data(chart_data)
                break

    def _update_time_series_charts(self, historical_df: pd.DataFrame):
        """Update all area charts in the presentation"""
        # Clean the DataFrame first
        historical_df = historical_df.replace([np.inf, -np.inf], 0)
        historical_df = historical_df.fillna(0)
        
        # Get last 24 months of data
        last_24_months = historical_df.tail(24)
        
        # Process each slide and its charts
        for slide_num, parties in PARTY_PAIRS.items():
            slide = self.prs.slides[slide_num - 1]
            
            # Get all charts in the slide
            charts = [shape.chart for shape in slide.shapes if shape.has_chart]
            
            # For slide 16 (CHP and AK Parti), reverse the order of parties
            if slide_num == 16:
                parties = list(reversed(parties))
            
            # Update each chart
            for i, party in enumerate(parties):
                if i < len(charts):
                    chart_data = CategoryChartData()
                    
                    # Use Date as categories with translation
                    chart_data.categories = self._translate_list(last_24_months['Months'].tolist(), 'dates')
                    
                    # Keep values as percentages (no division by 100)
                    if party == 'Kararsız':
                        values = last_24_months['Kararsız'].tolist()
                    else:
                        values = last_24_months[party].tolist()
                    
                    # Convert to float and handle any remaining NaN/INF
                    values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                    chart_data.add_series(self._translate(party, 'parties'), values)
                    charts[i].replace_data(chart_data)

    def _update_education_charts(self, historical_data: Dict[str, pd.DataFrame]):
        """Update education breakdown time series charts"""
        party_sheet_mapping = {
            'AK Parti': 'party_votes_education_akp',
            'CHP': 'party_votes_education_chp',
            'DEM Parti': 'party_votes_education_dem',
            'İYİ Parti': 'party_votes_education_iyip',
            'MHP': 'party_votes_education_mhp',
            'Kararsız': 'party_votes_education_kararsiz',
            'Oy Kullanmam': 'party_votes_education_absent'
        }
        
        for party, sheet_name in party_sheet_mapping.items():
            if sheet_name not in historical_data:
                continue
                
            df = historical_data[sheet_name]
            
            # Find the chart by name
            chart_name = f'education_{sheet_name.split("_")[-1]}'
            for slide in self.prs.slides:
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == chart_name), None)
                if chart:
                    chart_data = CategoryChartData()
                    
                    # Use dates from Months column with translation
                    chart_data.categories = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    
                    # Add each education level as a series with translations
                    for level in ['İlköğretim ve altı', 'Lise', 'Yüksekokul ve üzeri']:
                        values = df[level].tail(24).tolist()
                        values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                        chart_data.add_series(self._translate(level, 'education'), values)
                    
                    chart.replace_data(chart_data)
                    break

    def _update_age_charts(self, historical_data: Dict[str, pd.DataFrame]):
        """Update age breakdown time series charts"""
        party_sheet_mapping = {
            'AK Parti': 'party_votes_age_akp',
            'CHP': 'party_votes_age_chp',
            'DEM Parti': 'party_votes_age_dem',
            'İYİ Parti': 'party_votes_age_iyip',
            'MHP': 'party_votes_age_mhp',
            'Kararsız': 'party_votes_age_kararsiz',
            'Oy Kullanmam': 'party_votes_age_absent'
        }
        
        for party, sheet_name in party_sheet_mapping.items():
            if sheet_name not in historical_data:
                continue
                
            df = historical_data[sheet_name]
            
            # Find the chart by name
            chart_name = f'age_{sheet_name.split("_")[-1]}'
            for slide in self.prs.slides:
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == chart_name), None)
                if chart:
                    chart_data = CategoryChartData()
                    
                    # Use dates from Months column with translation
                    chart_data.categories = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    
                    # Add each age group as a series with translations
                    for group in ['18-34', '35-54', '55 ve üstü']:
                        values = df[group].tail(24).tolist()
                        values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                        chart_data.add_series(self._translate(group, 'age'), values)
                    
                    chart.replace_data(chart_data)
                    break

    def _update_2023_party_chart(self, retention_df: pd.DataFrame):
        """Update 2023 party retention time series chart"""
        # Find the chart by name
        chart = None
        for slide in self.prs.slides:
            chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == '2023_party'), None)
            if chart:
                break
        
        if chart:
            chart_data = CategoryChartData()
            
            # Use dates from Months column with translation
            chart_data.categories = self._translate_list(retention_df['Months'].tail(24).tolist(), 'dates')
            
            # Add series for each party with translations
            parties = ['AK Parti', 'CHP', 'DEM Parti', 'İYİ Parti', 'MHP']
            for party in parties:
                values = retention_df[party].tail(24).tolist()
                # Convert to float and handle any remaining NaN/INF
                values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                chart_data.add_series(self._translate(party, 'parties'), values)
            
            chart.replace_data(chart_data)

    def _update_econ_charts(self, historical_data: Dict[str, pd.DataFrame]):
        """Update all economic charts"""
        # Update main economic situation chart (Slide 32)
        if 'econ_main' in historical_data:
            df = historical_data['econ_main']
            for slide in self.prs.slides:
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'econ_main'), None)
                if chart:
                    chart_data = CategoryChartData()
                    chart_data.categories = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    
                    for response in ['Çok kötü / Kötü', 'Ne iyi ne kötü', 'Çok İyi / İyi']:
                        values = df[response].tail(24).tolist()
                        values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                        chart_data.add_series(self._translate(response, 'economy'), values)
                    
                    chart.replace_data(chart_data)
        
        # Update economic negative by party chart (Slide 33)
        if 'econ_negative_party' in historical_data:
            df = historical_data['econ_negative_party']
            for slide in self.prs.slides:
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'econ_negative_party'), None)
                if chart:
                    chart_data = CategoryChartData()
                    chart_data.categories = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    
                    for party in ['AK Parti', 'CHP', 'DEM Parti', 'İYİ Parti', 'MHP']:
                        values = df[party].tail(24).tolist()
                        values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                        chart_data.add_series(self._translate(party, 'parties'), values)
                    
                    chart.replace_data(chart_data)
        
        # Update economic negative by age chart (Slide 35)
        if 'econ_negative_age' in historical_data:
            df = historical_data['econ_negative_age']
            for slide in self.prs.slides:
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'econ_negative_age'), None)
                if chart:
                    chart_data = CategoryChartData()
                    chart_data.categories = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    
                    for age_group in ['18-34', '35-54', '55 ve üstü']:
                        values = df[age_group].tail(24).tolist()
                        values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                        chart_data.add_series(self._translate(age_group, 'age'), values)
                    
                    chart.replace_data(chart_data)
        
        # Update economic negative by education chart (Slide 37)
        if 'econ_negative_education' in historical_data:
            df = historical_data['econ_negative_education']
            for slide in self.prs.slides:
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'econ_negative_education'), None)
                if chart:
                    chart_data = CategoryChartData()
                    chart_data.categories = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    
                    for edu_level in ['İlköğretim ve altı', 'Lise', 'Yüksekokul ve üzeri']:
                        values = df[edu_level].tail(24).tolist()
                        values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                        chart_data.add_series(self._translate(edu_level, 'education'), values)
                    
                    chart.replace_data(chart_data)

        # Update main future economic situation chart (Slide 40)
        if 'econ_future_main' in historical_data:
            df = historical_data['econ_future_main']
            for slide in self.prs.slides:
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'econ_future_main'), None)
                if chart:
                    chart_data = CategoryChartData()
                    chart_data.categories = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    
                    for response in ['Çok Daha Kötü/Daha Kötü', 'Değişmez', 'Çok Daha İyi/Daha İyi']:
                        values = df[response].tail(24).tolist()
                        values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                        chart_data.add_series(self._translate(response, 'economy'), values)
                    
                    chart.replace_data(chart_data)
        
        # Update future economic negative by party chart (Slide 41)
        if 'econ_future_party' in historical_data:
            df = historical_data['econ_future_party']
            for slide in self.prs.slides:
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'econ_future_party'), None)
                if chart:
                    chart_data = CategoryChartData()
                    chart_data.categories = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    
                    for party in ['AK Parti', 'CHP', 'DEM Parti', 'İYİ Parti', 'MHP']:
                        values = df[party].tail(24).tolist()
                        values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                        chart_data.add_series(self._translate(party, 'parties'), values)
                    
                    chart.replace_data(chart_data)
        
        # Update future economic negative by age chart (Slide 43)
        if 'econ_future_age' in historical_data:
            df = historical_data['econ_future_age']
            for slide in self.prs.slides:
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'econ_future_age'), None)
                if chart:
                    chart_data = CategoryChartData()
                    chart_data.categories = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    
                    for age_group in ['18-34', '35-54', '55 ve üstü']:
                        values = df[age_group].tail(24).tolist()
                        values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                        chart_data.add_series(self._translate(age_group, 'age'), values)
                    
                    chart.replace_data(chart_data)

    def _update_politician_success_charts(self, historical_data: Dict[str, pd.DataFrame], current_success_data: pd.DataFrame):
        """Update all politician success charts"""
        print(f"\nUpdating politician success charts for language: {self.language}")
        print(f"Using presentation file: {self.output_path}")
        
        # Update current month's success rates (Slide 29)
        print("\nLooking for 'politician_success' chart...")
        chart_found = False
        for slide_num, slide in enumerate(self.prs.slides, 1):
            chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'politician_success'), None)
            if chart:
                print(f"Found 'politician_success' chart in slide {slide_num}")
                chart_found = True
                chart_data = CategoryChartData()
                # Create a copy of the DataFrame to avoid modifying the original
                current_data = current_success_data.copy()
                print(f"Original politician names: {current_data['Politician'].tolist()}")
                # Translate politician names
                current_data['Politician'] = current_data['Politician'].apply(lambda x: self._translate(x, 'politicians'))
                # Sort by Success Rate in descending order
                current_data = current_data.sort_values('Success Rate', ascending=False)
                print(f"Translated and sorted politician names: {current_data['Politician'].tolist()}")
                chart_data.categories = current_data['Politician'].tolist()
                chart_data.add_series(self._translate('Success Rate', 'chart_titles'), current_data['Success Rate'].tolist())
                try:
                    print("Attempting to replace chart data...")
                    chart.replace_data(chart_data)
                    print("Successfully replaced chart data")
                except Exception as e:
                    print(f"Error replacing chart data: {str(e)}")
                break
        
        if not chart_found:
            print("WARNING: 'politician_success' chart not found!")
        
        # Update main politicians historical chart (Slide 30)
        print("\nLooking for 'politician_success_main' chart...")
        if 'politician_success_main' in historical_data:
            df = historical_data['politician_success_main']
            print(f"Available columns in main data: {df.columns.tolist()}")
            chart_found = False
            for slide_num, slide in enumerate(self.prs.slides, 1):
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'politician_success_main'), None)
                if chart:
                    print(f"Found 'politician_success_main' chart in slide {slide_num}")
                    chart_found = True
                    chart_data = CategoryChartData()
                    dates = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    print(f"Translated dates: {dates}")
                    chart_data.categories = dates
                    
                    for politician in ['Recep Tayyip Erdoğan', 'Özgür Özel', 'Devlet Bahçeli', 
                                     'Ekrem İmamoğlu', 'Mansur Yavaş', 'Fatih Erbakan']:
                        if politician in df.columns:
                            print(f"Processing politician: {politician}")
                            values = df[politician].tail(24).tolist()
                            # Convert zeros to None instead of 0
                            values = [None if pd.isna(x) or not np.isfinite(x) or x == 0 else float(x) for x in values]
                            translated_name = self._translate(politician, 'politicians')
                            print(f"Translated name: {translated_name}")
                            chart_data.add_series(translated_name, values)
                        else:
                            print(f"WARNING: Politician {politician} not found in data!")
                    
                    try:
                        print("Attempting to replace chart data...")
                        chart.replace_data(chart_data)
                        print("Successfully replaced chart data")
                    except Exception as e:
                        print(f"Error replacing chart data: {str(e)}")
                    break
            
            if not chart_found:
                print("WARNING: 'politician_success_main' chart not found!")
        else:
            print("WARNING: 'politician_success_main' data not found in historical data!")
        
        # Update secondary politicians historical chart (Slide 30)
        print("\nLooking for 'politician_success_second' chart...")
        if 'politician_success_second' in historical_data:
            df = historical_data['politician_success_second']
            print(f"Available columns in secondary data: {df.columns.tolist()}")
            chart_found = False
            for slide_num, slide in enumerate(self.prs.slides, 1):
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'politician_success_second'), None)
                if chart:
                    print(f"Found 'politician_success_second' chart in slide {slide_num}")
                    chart_found = True
                    chart_data = CategoryChartData()
                    dates = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    print(f"Translated dates: {dates}")
                    chart_data.categories = dates
                    
                    for politician in ['Muharrem İnce', 'Erkan Baş', 'Ümit Özdağ', 'Müsavat Dervişoğlu',
                                     'Tülay Hatimoğulları Oruç', 'Yavuz Ağıralioğlu', 'Mahmut Arıkan']:
                        if politician in df.columns:
                            print(f"Processing politician: {politician}")
                            values = df[politician].tail(24).tolist()
                            # Convert zeros to None instead of 0
                            values = [None if pd.isna(x) or not np.isfinite(x) or x == 0 else float(x) for x in values]
                            translated_name = self._translate(politician, 'politicians')
                            print(f"Translated name: {translated_name}")
                            chart_data.add_series(translated_name, values)
                        else:
                            print(f"WARNING: Politician {politician} not found in data!")
                    
                    try:
                        print("Attempting to replace chart data...")
                        chart.replace_data(chart_data)
                        print("Successfully replaced chart data")
                    except Exception as e:
                        print(f"Error replacing chart data: {str(e)}")
                    break
            
            if not chart_found:
                print("WARNING: 'politician_success_second' chart not found!")
        else:
            print("WARNING: 'politician_success_second' data not found in historical data!")

    def _update_subsistence_charts(self, historical_data: Dict[str, pd.DataFrame]):
        """Update all subsistence charts"""
        # Update main subsistence chart (Slide 49)
        if 'subsistence' in historical_data:
            df = historical_data['subsistence']
            for slide in self.prs.slides:
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'subsistence'), None)
                if chart:
                    chart_data = CategoryChartData()
                    chart_data.categories = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    
                    responses = [
                        'Gelirim giderimi karşılamadı.',
                        'Gelirim giderimi ucu ucuna karşıladı.',
                        'Gelirim giderlerimin üzerinde oldu.',
                        'Gelirim giderlerimi fazlasıyla karşıladı.'
                    ]
                    
                    for response in responses:
                        values = df[response].tail(24).tolist()
                        values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                        chart_data.add_series(self._translate(response, 'subsistence'), values)
                    
                    chart.replace_data(chart_data)
        
        # Update subsistence by party chart (Slide 51)
        if 'subsistence_party' in historical_data:
            df = historical_data['subsistence_party']
            for slide in self.prs.slides:
                chart = next((shape.chart for shape in slide.shapes if shape.has_chart and shape.name == 'subsistence_party'), None)
                if chart:
                    chart_data = CategoryChartData()
                    chart_data.categories = self._translate_list(df['Months'].tail(24).tolist(), 'dates')
                    
                    for party in ['AK Parti', 'CHP', 'DEM Parti', 'İYİ Parti', 'MHP']:
                        values = df[party].tail(24).tolist()
                        values = [float(0) if pd.isna(x) or not np.isfinite(x) else float(x) for x in values]
                        chart_data.add_series(self._translate(party, 'parties'), values)
                    
                    chart.replace_data(chart_data)

    def _prepare_sorted_data(self, party_data: pd.DataFrame) -> pd.DataFrame:
        """Prepare sorted data with 'Diğer' always at the bottom"""
        diger_data = party_data[party_data['Party'] == 'Diğer']
        other_parties = party_data[party_data['Party'] != 'Diğer']
        sorted_parties = other_parties.sort_values('Percentage', ascending=False)
        if not diger_data.empty:
            sorted_parties = pd.concat([sorted_parties, diger_data])
        return sorted_parties